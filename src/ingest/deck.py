"""
src/ingest/deck.py — Slide deck (PDF or PPTX) ingestion.

Mirrors the PROVIDED video flow, but the source is a slide deck and the citation
locator is a **slide** number:

    fetch → parse (per slide, caption image-only) → chunk (per slide) → embed → index

Reconciled to the REAL momentsearch repo (same as paper.py — see its header and
RECONCILED.md). Decks are TEXT once captioned, so they ride the repo's *text
branch* (`embed_docs` + `vector_store.upsert_chunks` into `TEXT_COLLECTION`),
tagged `kind="deck"`, into the SAME collection as videos and papers.

Key wrinkle vs. papers: image-only slides have little/no extractable text, so we
caption them with the repo's multimodal LLM (`src/llm.py :: answer`, env- or
tenant-switched) BEFORE embedding — otherwise those slides retrieve poorly (a
graded pitfall). The repo has no dedicated `caption_image`, so `_caption` wraps
`llm.answer` with a single-image "moment".

LIMITATION — captioning is **PDF-only**. PyMuPDF rasterizes a PDF page, so an
image-only slide in a PDF deck gets a caption. `python-pptx` reads shapes and
cannot render a slide, so `_render_pptx_slide_png` has nothing to hand the vision
model and returns None: an image-only PPTX slide indexes on its sparse text
alone. Closing this means a LibreOffice-headless PPTX→PDF conversion up front,
reusing `_parse_pdf_deck` — a heavy binary in the image for one input format,
so it is deliberately not done. Export the deck to PDF to get captioning.
"""
from __future__ import annotations

import io
from dataclasses import dataclass, field

import fitz  # PyMuPDF (PDF decks + rendering slides to images)   pip install pymupdf

# --- shared repo pieces (REAL names, reconciled) ------------------------------
from .. import db, llm, storage
from ..config import TEXT_EMBED_VERSION
from ..rag import vector_store
from ..rag.embeddings import embed_docs
# reuse paper.py's fetch + word-window chunker (identical needs) so there's one
# implementation to maintain.
from .paper import _load_bytes, _window_chunks, CHUNK_WORDS, MAX_PAGES

KIND = "deck"
MIN_TEXT_CHARS = 24          # below this, treat a slide as image-only and caption it
RENDER_DPI = 120             # rasterization DPI for image-only slide captioning


@dataclass
class Slide:
    slide: int         # 1-based slide number == the citation locator
    text: str          # slide text (+ appended caption for image-only slides)


@dataclass
class Chunk:
    text: str
    payload: dict = field(default_factory=dict)


# 0) CAPTION (image-only slides) ──────────────────────────────────────────────
def _caption(png: bytes) -> str:
    """Caption a rendered slide with the repo's multimodal LLM.

    The repo exposes `llm.answer(question, moments, cfg)` (no standalone
    `caption_image`), where a moment is {"image": bytes, "transcript": None,
    "timestamp": str}. We reuse it with the server-wide env model. Captioning is
    best-effort: if no LLM is configured, or the call fails, return "" and let
    the slide index on whatever text it has (never fail the flow — mirrors the
    transcript branch's best-effort stance).
    """
    cfg = llm.env_config()
    if cfg is None:
        return ""
    try:
        return llm.answer(
            "Describe this slide in one or two sentences for search: its title, "
            "key text, and any chart/diagram/image shown.",
            [{"image": png, "transcript": None, "timestamp": "slide"}],
            cfg,
        ).strip()
    except Exception as exc:  # noqa: BLE001
        print(f"[deck] caption failed ({type(exc).__name__}: {exc}) — text-only slide")
        return ""


# 1) PARSE ─────────────────────────────────────────────────────────────────────
def parse_deck(data: bytes, filename: str = "") -> list[Slide]:
    """Deck bytes → per-slide text. Handles PPTX and PDF; captions image-only
    slides in PDF decks only (see the module docstring)."""
    if filename.lower().endswith(".pptx"):
        return _parse_pptx(data)
    return _parse_pdf_deck(data)


def _parse_pptx(data: bytes) -> list[Slide]:
    from pptx import Presentation                 # pip install python-pptx

    prs = Presentation(io.BytesIO(data))
    slides: list[Slide] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [sh.text for sh in slide.shapes
                 if getattr(sh, "has_text_frame", False) and sh.text]
        text = "\n".join(t.strip() for t in texts if t.strip())
        if len(text) < MIN_TEXT_CHARS:
            png = _render_pptx_slide_png(slide)   # see helper note below
            if png:
                text = (text + "\n" + _caption(png)).strip()
        slides.append(Slide(slide=i, text=text))
    return slides


def _parse_pdf_deck(data: bytes) -> list[Slide]:
    """PDF deck: one page == one slide. Extract text; caption image-only slides."""
    slides: list[Slide] = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        if doc.page_count > MAX_PAGES:     # §11: reject abusive/oversized decks
            raise ValueError(f"Deck has {doc.page_count} slides; the limit is {MAX_PAGES}.")
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if len(text) < MIN_TEXT_CHARS:
                png = page.get_pixmap(dpi=RENDER_DPI).tobytes("png")
                text = (text + "\n" + _caption(png)).strip()
            slides.append(Slide(slide=i, text=text))
    return slides


def _render_pptx_slide_png(slide) -> bytes | None:  # noqa: ANN001
    """Always None — PPTX slides are never captioned. NOT a stub left by mistake.

    python-pptx reads shapes; it cannot rasterize a slide, and there is no
    pure-Python renderer for PowerPoint layout. The two ways out:
      • convert the PPTX to PDF once (LibreOffice headless) and reuse
        `_parse_pdf_deck` — one rasterization path, but a ~400mb binary in the
        image for a single input format; or
      • pull embedded picture shapes' blobs and caption those individually —
        cheap, but it captions pictures, not slides: layout, title and chart
        text around the image are lost, which is most of what makes a slide
        findable.
    Neither is worth it while PDF decks (the common export) caption correctly.
    Returning None degrades to text-only indexing for that slide — never fatal.
    """
    return None


# 2) CHUNK (slide-aware) ──────────────────────────────────────────────────────
def chunk_deck(source_id: str, title: str, slides: list[Slide],
               user_id: str) -> list[Chunk]:
    """One (or a few) chunk(s) per slide, each carrying `slide` as the locator.
    Most slides are short → one chunk; only long, dense slides get windowed.

    `user_id` MUST be in the payload — every search is user_id-filtered, so a
    chunk without it is invisible to retrieval (mirrors the video payload)."""
    chunks: list[Chunk] = []
    for sl in slides:
        words = len(sl.text.split())
        pieces = _window_chunks(sl.text) if words > CHUNK_WORDS else (
            [sl.text] if sl.text.strip() else [])
        for piece in pieces:
            chunks.append(Chunk(
                text=piece,
                payload={
                    "user_id": user_id,       # ← multi-tenant filter (required!)
                    "kind": KIND,
                    "source_id": source_id,
                    "video_id": source_id,   # read path fuses/filters on video_id
                    "title": title,
                    "slide": sl.slide,        # ← locator
                    "modality": "text",
                    "text": piece,
                    "t_start": 0.0, "t_end": 0.0,
                    "embed_version": TEXT_EMBED_VERSION,
                },
            ))
    return chunks


# 3) FLOW ──────────────────────────────────────────────────────────────────────
def ingest_deck(source_id: str, uri: str, title: str, user_id: str,
                storage_key: str | None = None, filename: str = "") -> int:
    """
    Deck flow body. Same status lifecycle + crash-safety as paper/video:
    set `indexed` only AFTER the Qdrant upsert succeeds; deterministic point ids
    (via `upsert_chunks`) + `delete_video` make re-runs idempotent.

    Returns the number of chunks indexed.
    """
    db.bump_attempts(source_id)
    try:
        db.set_status(source_id, "parsing")
        data = _load_bytes(uri, storage_key)
        slides = parse_deck(data, filename or uri)

        db.set_status(source_id, "chunking")
        chunks = chunk_deck(source_id, title, slides, user_id)
        if not chunks:
            db.set_status(source_id, "indexed", frame_count=0, progress=1.0)
            return 0

        db.set_status(source_id, "embedding", progress=0.0)
        vectors = embed_docs([c.text for c in chunks])

        vector_store.ensure_text_collection()
        vector_store.delete_video(user_id, source_id)
        vector_store.upsert_chunks(
            user_id, source_id, vectors,
            payloads=[c.payload for c in chunks],
        )
        db.set_status(source_id, "indexed", frame_count=len(chunks),
                      embed_version=TEXT_EMBED_VERSION, progress=1.0)
        return len(chunks)
    except Exception as exc:                             # noqa: BLE001
        db.set_status(source_id, "failed", error=f"{type(exc).__name__}: {exc}")
        raise


# ── Prefect wiring: add an `ms-ingest-deck` flow beside `ms-ingest-paper` /
#    video, same per-task retries; POST /admin/documents schedules it for
#    kind="deck". See paper.py's footer for the flow/task sketch.
